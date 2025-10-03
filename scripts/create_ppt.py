from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Challenge III: Automotive — Agentic AI for Predictive Maintenance"
    subtitle.text = (
        "Master-Agent orchestration for proactive service scheduling, persuasive voice engagement,\n"
        "RCA/CAPA feedback to manufacturing, and UEBA-driven security"
    )
    return slide


def add_section_header(slide, text):
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9.0)
    height = Inches(0.6)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(34, 65, 121)  # Deep blue
    line = shape.line
    line.color.rgb = RGBColor(34, 65, 121)
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def add_bulleted_textbox(slide, title_text, bullets, left=0.5, top=1.3, width=4.5, height=4.5):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(18)
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1
        p.font.size = Pt(14)
    return shape


def add_caption_box(slide, caption, left=0.5, top=6.0, width=9.0, height=0.6):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
    p.alignment = PP_ALIGN.CENTER


def add_architecture_block(slide, title, desc, left, top, width=2.8, height=1.3, color=(223, 231, 243)):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*color)
    rect.line.color.rgb = RGBColor(34, 65, 121)
    tf = rect.text_frame
    tf.text = title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(14)
    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(12)
    return rect


def add_connector(slide, x1, y1, x2, y2):
    # Adds a straight connector line between points (inches)
    line = slide.shapes.add_line(Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = RGBColor(34, 65, 121)
    line.line.width = Pt(2)
    return line


def slide_monitoring_and_prediction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_section_header(slide, "1) Continuous Monitoring and Predictive Failure Detection")

    add_bulleted_textbox(
        slide,
        "Data + Pipeline",
        bullets=[
            "Telematics API (mock): real-time engine, braking, battery, temp, vibration",
            "Historical maintenance logs + DTCs; CAPA/RCA records",
            "Feature engineering: trend slopes, anomaly scores, usage intensity, dwell time",
            "Models: gradient boosting + survival analysis for time-to-failure",
            "Priority scoring: severity × likelihood × time-to-failure",
        ],
        left=0.5,
        top=1.3,
        width=4.5,
        height=3.6,
    )

    add_bulleted_textbox(
        slide,
        "Master-Agent Orchestration",
        bullets=[
            "Data Analysis Agent: streaming inference + early warnings",
            "Diagnosis Agent: component-level failure probabilities and priorities",
            "Master Agent: consolidates alerts, enforces UEBA policy checks",
            "Edge cases: urgent alerts, multi-vehicle fleets, holiday loads",
        ],
        left=5.2,
        top=1.3,
        width=4.3,
        height=3.6,
    )

    add_bulleted_textbox(
        slide,
        "Illustrative Output",
        bullets=[
            "VHC-07: Brake pad failure risk 0.82, ETA 11 days — priority High",
            "VHC-03: Battery health degradation, replacing cell pack — priority Medium",
            "VHC-10: Engine coolant temp anomalies under hill loads — priority High",
        ],
        left=0.5,
        top=5.0,
        width=9.0,
        height=1.6,
    )

    add_caption_box(
        slide,
        "Streaming telemetry + historical logs feed models; Master Agent coordinates diagnosis and flags risks early.",
    )


def slide_demand_forecasting_and_scheduling(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "2) Forecasting Service Demand and Autonomous Scheduling")

    add_bulleted_textbox(
        slide,
        "Service Demand Forecasting",
        bullets=[
            "Aggregate patterns: mileage, duty cycles, climate, route topology",
            "Seasonality: pre-festive travel spikes; monsoon wear patterns",
            "Workload prediction per service center to optimize staffing and bays",
        ],
        left=0.5,
        top=1.3,
        width=4.5,
        height=2.8,
    )

    add_bulleted_textbox(
        slide,
        "Scheduling Agent (Autonomous)",
        bullets=[
            "Checks capacity via Scheduler API; proposes optimal slots",
            "Balances customer preferences, travel distance, urgency",
            "Automatically books, reschedules, or escalates urgent cases",
        ],
        left=5.2,
        top=1.3,
        width=4.3,
        height=2.8,
    )

    add_bulleted_textbox(
        slide,
        "Persuasive Voice Agent Snippet",
        bullets=[
            "“Hi Ananya, your vehicle reports increased brake wear with high confidence.",
            "To avoid downtime, shall I book a 10:30 AM slot tomorrow at Indiranagar?",
            "We’ll perform predictive checks and fast-track service. Does that work?”",
        ],
        left=0.5,
        top=4.4,
        width=9.0,
        height=2.2,
    )

    add_caption_box(
        slide,
        "Demand forecasts inform proactive slot allocation; Voice agent drives conversion; app notifications backstop reminders.",
    )


def slide_customer_engagement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "3) Persuasive Customer Engagement via Voice Agent")

    add_bulleted_textbox(
        slide,
        "Conversational Strategy",
        bullets=[
            "Explain condition in plain language; quantify risk and benefit",
            "Offer choices: nearest centers, earliest slots, shuttle/valet options",
            "Handle declines: suggest alternative dates, reinforce safety/economics",
            "Confirm via OTP/App; send checklist and ETA",
        ],
        left=0.5,
        top=1.3,
        width=4.5,
        height=3.1,
    )

    add_bulleted_textbox(
        slide,
        "Sample Dialogs",
        bullets=[
            "Owner: “Busy this week.” Agent: “Understood. Next Tue 7:30 AM or Thu 6 PM?”",
            "Owner: “What if I wait?” Agent: “Risk of brake fade is high under city traffic.",
            "A quick service prevents rotor damage and saves ₹4,000–₹7,000.”",
        ],
        left=5.2,
        top=1.3,
        width=4.3,
        height=3.1,
    )

    add_bulleted_textbox(
        slide,
        "Tracking and Follow-up",
        bullets=[
            "Feedback Agent: post-service CSAT, captures notes and parts replaced",
            "Updates maintenance records; closes loop with Master Agent",
            "Reinforcement signals for model recalibration",
        ],
        left=0.5,
        top=4.8,
        width=9.0,
        height=1.8,
    )

    add_caption_box(
        slide,
        "Voice-first experience increases booking rates; empathetic, informative prompts reduce deferment.",
    )


def slide_rca_capa_feedback(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "4) RCA/CAPA Insights and Manufacturing Feedback Loop")

    add_bulleted_textbox(
        slide,
        "Manufacturing Quality Insights Module",
        bullets=[
            "Correlate predicted failures with historical CAPA/RCA and part lots",
            "Detect recurring defects by geography/usage and supplier batches",
            "Recommend design change, process control, or supplier corrective action",
        ],
        left=0.5,
        top=1.3,
        width=4.5,
        height=2.6,
    )

    add_bulleted_textbox(
        slide,
        "Example Insight",
        bullets=[
            "Brake pad compound batch B47 shows accelerated wear in humid metros",
            "CAPA: switch to revised resin; update QC humidity soak test protocol",
            "Service advisory: shorten inspection interval in monsoon months",
        ],
        left=5.2,
        top=1.3,
        width=4.3,
        height=2.6,
    )

    # Simple architecture diagram blocks
    add_architecture_block(
        slide, "After-sales Data", "Telemetry + Maintenance + DTCs", left=0.5, top=3.9
    )
    add_architecture_block(
        slide, "Insights Engine", "RCA/CAPA correlation + recurrence detection", left=3.4, top=3.9
    )
    add_architecture_block(
        slide, "Manufacturing", "Design/process updates + supplier actions", left=6.3, top=3.9
    )
    add_connector(slide, 3.3, 4.55, 3.4, 4.55)
    add_connector(slide, 6.2, 4.55, 6.3, 4.55)

    add_bulleted_textbox(
        slide,
        "Closed-loop Outcomes",
        bullets=[
            "Reduced repeat defects; faster CAPA cycles",
            "Fewer breakdowns; lower warranty spend",
            "Data-driven design improvements",
        ],
        left=0.5,
        top=5.6,
        width=9.0,
        height=1.4,
    )

    add_caption_box(
        slide,
        "Predicted failures + RCA/CAPA patterns inform manufacturing changes; feedback reduces recurring defects.",
    )


def slide_ueba(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "5) UEBA in Action — Securing Autonomous Agent Orchestration")

    add_bulleted_textbox(
        slide,
        "UEBA Baselines",
        bullets=[
            "Master Agent: allowed orchestration + audit log writes",
            "Data Analysis Agent: telematics read; no booking writes",
            "Scheduling Agent: scheduler API write; no telematics read",
            "Customer Agent: comms channels only; no manufacturing DB access",
        ],
        left=0.5,
        top=1.3,
        width=4.5,
        height=3.0,
    )

    add_bulleted_textbox(
        slide,
        "Anomaly Examples",
        bullets=[
            "Scheduling Agent attempts telematics read → alert + block",
            "Customer Agent tries to modify slots → quarantine conversation",
            "Unexpected workflow jump (booking without diagnosis) → require approval",
        ],
        left=5.2,
        top=1.3,
        width=4.3,
        height=3.0,
    )

    add_bulleted_textbox(
        slide,
        "Response Flow",
        bullets=[
            "Detect → alert → auto-block → escalate to SOC → review/audit",
            "Continuous learning: update baselines and policies",
            "Compliance: full traceability of agent actions",
        ],
        left=0.5,
        top=4.6,
        width=9.0,
        height=1.8,
    )

    add_caption_box(
        slide,
        "UEBA prevents unauthorized actions and anomalous agent behavior while preserving autonomy.",
    )


def build_presentation(output_path: str):
    prs = Presentation()

    # Slide 1 — Title/Hero
    add_title_slide(prs)

    # Slide 2 — Monitoring & Prediction
    slide_monitoring_and_prediction(prs)

    # Slide 3 — Demand Forecasting & Scheduling with Voice Agent
    slide_demand_forecasting_and_scheduling(prs)

    # Slide 4 — Customer Engagement (Persuasive Voice)
    slide_customer_engagement(prs)

    # Slide 5 — RCA/CAPA + Manufacturing Feedback
    slide_rca_capa_feedback(prs)

    # Slide 6 — UEBA in action
    slide_ueba(prs)

    prs.save(output_path)


if __name__ == "__main__":
    # Default output path
    build_presentation("out/Challenge_III_Automotive_Agentic_AI.pptx")